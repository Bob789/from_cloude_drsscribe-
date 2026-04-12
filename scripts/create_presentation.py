"""
Generate a professional portfolio PowerPoint presentation
for the Doctor-Scribe ecosystem (3 products).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml.etree import SubElement

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)
BG_CARD    = RGBColor(0x16, 0x21, 0x3B)
ACCENT     = RGBColor(0x38, 0xBD, 0xF8)  # bright blue
ACCENT2    = RGBColor(0x22, 0xD3, 0xEE)  # cyan
GREEN      = RGBColor(0x4A, 0xDE, 0x80)
ORANGE     = RGBColor(0xFB, 0x92, 0x3C)
PURPLE     = RGBColor(0xA7, 0x8B, 0xFA)
RED        = RGBColor(0xF8, 0x71, 0x71)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
LIGHT      = RGBColor(0xE2, 0xE8, 0xF0)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE,
                 bold=False, alignment=PP_ALIGN.RIGHT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=14, color=LIGHT,
                    bullet_color=ACCENT, alignment=PP_ALIGN.RIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.alignment = alignment
        p.space_after = Pt(6)
    return txBox


def add_icon_card(slide, left, top, width, height, icon, title, desc, accent=ACCENT):
    card = add_shape_rect(slide, left, top, width, height, BG_CARD, accent)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.15), width - Inches(0.3), Inches(0.5),
                 icon, font_size=28, color=accent, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.7), width - Inches(0.3), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(1.05), width - Inches(0.2), Inches(0.8),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    return card


def add_tooltip(shape, tooltip_text):
    """Add ScreenTip (hover tooltip) to a shape — visible in PowerPoint presentation mode."""
    sp = shape._element
    nvSpPr = sp.find(qn('p:nvSpPr'))
    if nvSpPr is None:
        return
    cNvPr = nvSpPr.find(qn('p:cNvPr'))
    if cNvPr is None:
        return
    hlinkClick = SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), '')
    hlinkClick.set('tooltip', tooltip_text)
    hlinkClick.set('action', 'ppaction://noaction')


def add_tooltip_overlay(slide, left, top, width, height, tooltip_text):
    """Add a transparent overlay shape with a hover tooltip covering the card area."""
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    overlay.fill.background()
    overlay.line.fill.background()
    add_tooltip(overlay, tooltip_text)
    return overlay


# ════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Decorative top bar
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.2),
             "MedScribe AI Ecosystem", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(0.8),
             "Full-Stack · Data Science · AI · Medical SaaS",
             font_size=26, color=ACCENT, bold=False, alignment=PP_ALIGN.CENTER)

# Three product cards
card_w = Inches(3.6)
card_h = Inches(2.2)
gap = Inches(0.5)
start_x = Inches(1.3)
y = Inches(4.2)

products = [
    ("🔬", "Data Science Platform", "חיזוי רפואי + השוואת מודלים\nStreamlit · FastAPI · scikit-learn", PURPLE),
    ("🏥", "Doctor Scribe AI", "ניהול מטופלים + תמלול + סיכום AI\nFastAPI · Flutter · Whisper · Claude", ACCENT),
    ("🌐", "Parent Website", "מאמרים + פורום + כתיבה אוטומטית\nNext.js · Tailwind · PostgreSQL", GREEN),
]

product_tooltips = [
    "פלטפורמת ML: העלאת CSV → הרצת 12+ מודלים → השוואה אוטומטית → חיזוי",
    "SaaS רפואי: הקלטת ביקור → תמלול Whisper → סיכום Claude → חיפוש ותיוג",
    "אתר שיווקי: 150+ מאמרים רפואיים · פורום SO-style · כתיבת תוכן AI · SEO",
]

for i, (icon, title, desc, color) in enumerate(products):
    x = start_x + i * (card_w + gap)
    add_icon_card(slide, x, y, card_w, card_h, icon, title, desc, color)
    add_tooltip_overlay(slide, x, y, card_w, card_h, product_tooltips[i])

add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5),
             "Yossi · Full-Stack & Data Science Developer · 2026",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "Architecture Overview — 9 Microservices", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.5),
             "Docker Compose · Nginx Reverse Proxy · SSL/TLS 1.3 · Health Checks",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# Service boxes grid
services = [
    ("Nginx", "Reverse Proxy + SSL\nRate Limiting · HSTS", ACCENT),
    ("FastAPI Backend", "Python 3.11 · 60+ Endpoints\n18 Routers · 21 Services", GREEN),
    ("Flutter Frontend", "Dart 3 · 15 Screens\nRiverpod · 8 Languages", ORANGE),
    ("Next.js Website", "React 18 · TypeScript\n17+ Routes · SEO", PURPLE),
    ("Celery Workers", "Async Tasks · STT\nLLM Summary · Tagging", RED),
    ("PostgreSQL 16", "16 DB Models\nAlembic Migrations", ACCENT2),
    ("Redis 7", "Cache · Sessions\nTask Queue (Broker)", ORANGE),
    ("MinIO / S3", "Audio Storage\nAES-256 Encryption", GREEN),
    ("Meilisearch", "Full-Text Search\nSemantic Indexing", PURPLE),
]

cols = 3
box_w = Inches(3.8)
box_h = Inches(1.4)
x_start = Inches(0.75)
y_start = Inches(1.7)
gap_x = Inches(0.35)
gap_y = Inches(0.2)

service_tooltips = [
    "Reverse Proxy: ניתוב תעבורה, SSL/TLS 1.3, Rate Limiting, HSTS",
    "שרת API בפייתון: 60+ endpoints, אימות JWT, 18 routers, 21 services",
    "אפליקציית Web: 15 מסכים, הקלטה, Riverpod, 8 שפות, RTL/LTR",
    "אתר שיווקי: React 18, TypeScript, 17+ routes, SSR, SEO",
    "עובדים אסינכרוניים: תמלול Whisper, סיכום Claude, תיוג אוטומטי",
    "DB ראשי: 16 מודלים, Alembic migrations, אחסון מוצפן",
    "מטמון: sessions, task queue ל-Celery, מאיץ תגובות API",
    "אחסון קבצים: אודיו מוצפן AES-256, signed URLs, chunked upload",
    "חיפוש: Full-text + Semantic, פילטור לפי תגיות ותאריכים",
]

for idx, (name, desc, color) in enumerate(services):
    col = idx % cols
    row = idx // cols
    x = x_start + col * (box_w + gap_x)
    y = y_start + row * (box_h + gap_y)
    card = add_shape_rect(slide, x, y, box_w, box_h, BG_CARD, color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.15), box_w - Inches(0.4), Inches(0.4),
                 name, font_size=16, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), box_w - Inches(0.4), Inches(0.7),
                 desc, font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)
    add_tooltip_overlay(slide, x, y, box_w, box_h, service_tooltips[idx])


# ════════════════════════════════════════════════════
# SLIDE 3 — Data Science Platform
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "🔬 Data Science Platform", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
             "חיזוי רפואי מבוסס Machine Learning + השוואה אוטומטית בין מודלים",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Left column — ML Models
add_shape_rect(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.2), BG_CARD, PURPLE)
add_text_box(slide, Inches(0.7), Inches(1.85), Inches(5.4), Inches(0.4),
             "🤖 12+ ML Models Supported", font_size=18, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

models_text = [
    "▸ Classification: Logistic Reg · Decision Tree · Random Forest",
    "▸ Classification: Gradient Boosting · SVM · Naive Bayes · KNN",
    "▸ Regression: Linear · Ridge/Lasso · RF · Gradient Boosting",
    "▸ Auto-Compare: מריץ את כל המודלים על אותו דאטה",
    "▸ Metrics: Accuracy, F1, ROC-AUC, RMSE, R², Confusion Matrix",
    "▸ Cross-Validation: 3-10 folds + Hold-out test set",
    "▸ Preprocessing: One-hot encoding + Mean imputation",
]
add_bullet_list(slide, Inches(0.7), Inches(2.35), Inches(5.4), Inches(3.8),
                models_text, font_size=13, color=LIGHT, alignment=PP_ALIGN.LEFT)

# Right column — Tech & Features
add_shape_rect(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.4), BG_CARD, ACCENT2)
add_text_box(slide, Inches(7.0), Inches(1.85), Inches(5.6), Inches(0.4),
             "⚡ Tech Stack", font_size=18, color=ACCENT2, bold=True, alignment=PP_ALIGN.CENTER)

tech_items = [
    "FastAPI + Uvicorn — REST API לאימון וחיזוי",
    "Streamlit — Dashboard אינטראקטיבי עם גרפים",
    "scikit-learn + pandas + NumPy — ML Pipeline",
    "PostgreSQL — אחסון מודלים + משתמשים",
    "Docker — Containerized deployment",
]
add_bullet_list(slide, Inches(7.0), Inches(2.35), Inches(5.6), Inches(1.6),
                tech_items, font_size=13, color=LIGHT, alignment=PP_ALIGN.LEFT)

add_shape_rect(slide, Inches(6.8), Inches(4.3), Inches(6.0), Inches(2.6), BG_CARD, ORANGE)
add_text_box(slide, Inches(7.0), Inches(4.45), Inches(5.6), Inches(0.4),
             "💡 Unique Value", font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

unique_items = [
    "Token System — מערכת מטבעות לשימוש (SaaS monetization)",
    "Dual Interface — API לאוטומציה + UI לאנליסט",
    "Auto Task Detection — זיהוי אוטומטי classification/regression",
    "3 Demo Datasets רפואיים מובנים (Heart, Housing, Smoking)",
    "Model Registry — שמירה, טעינה ומחיקה של מודלים",
]
add_bullet_list(slide, Inches(7.0), Inches(4.95), Inches(5.6), Inches(1.8),
                unique_items, font_size=13, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 4 — DS Platform — Model Comparison Flow
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "DS Platform — Model Comparison Pipeline", font_size=36, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

# Pipeline steps
steps = [
    ("📁", "Upload CSV", "משתמש מעלה\nדאטהסט CSV", ACCENT2),
    ("🔍", "Auto-Detect", "זיהוי Classification\nvs Regression", PURPLE),
    ("⚙️", "Preprocess", "One-Hot Encoding\nMean Imputation", GREEN),
    ("🤖", "Train All", "12+ מודלים\nרצים במקביל", ORANGE),
    ("📊", "Evaluate", "CV + Hold-out\nMetrics מלאים", ACCENT),
    ("🏆", "Compare", "דירוג מודלים\nלפי כל מטריקה", RED),
]

step_w = Inches(1.8)
step_h = Inches(2.8)
x_start = Inches(0.5)
gap = Inches(0.27)
y = Inches(1.5)

pipeline_tooltips = [
    "משתמש מעלה קובץ CSV עם נתונים, המערכת מזהה עמודות ויעד",
    "זיהוי אוטומטי אם המשימה Classification או Regression",
    "ניקוי: One-Hot Encoding לקטגוריות, Mean Imputation לחסרים",
    "הרצת 12+ מודלים במקביל על אותו דאטה עם Cross-Validation",
    "חישוב מטריקות: Accuracy, F1, ROC-AUC, RMSE, R², Confusion Matrix",
    "דירוג כל המודלים לפי כל מטריקה, בחירת המודל הטוב ביותר",
]

for i, (icon, title, desc, color) in enumerate(steps):
    x = x_start + i * (step_w + gap)
    add_shape_rect(slide, x, y, step_w, step_h, BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.2), step_w, Inches(0.5),
                 icon, font_size=32, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.8), step_w - Inches(0.2), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.3), step_w - Inches(0.2), Inches(1.2),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_tooltip_overlay(slide, x, y, step_w, step_h, pipeline_tooltips[i])
    # Arrow (except last)
    if i < len(steps) - 1:
        add_text_box(slide, x + step_w, y + Inches(1.0), Inches(gap), Inches(0.5),
                     "→", font_size=24, color=GRAY, alignment=PP_ALIGN.CENTER)

# Business value box
add_shape_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.3), BG_CARD, GREEN)
add_text_box(slide, Inches(0.7), Inches(4.85), Inches(11.9), Inches(0.4),
             "📈 Business Value", font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

biz_items = [
    "▸ רופאים/חוקרים יכולים להריץ חיזוי רפואי ללא ידע בתכנות — רק מעלים CSV ומקבלים תוצאות",
    "▸ השוואה אוטומטית חוסכת שעות של עבודת Data Scientist — במקום לבדוק מודל-מודל, הכל רץ בקליק אחד",
    "▸ Token System מאפשר מודל SaaS — תשלום לפי שימוש, scalable revenue",
    "▸ 3 דמו-דאטהסטים רפואיים מובנים — מאפשרים הדגמה מיידית ללקוחות פוטנציאליים",
]
add_bullet_list(slide, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.5),
                biz_items, font_size=14, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 5 — Doctor Scribe AI
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "🏥 Doctor Scribe AI — Medical Visit Management", font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
             "הקלטה · תמלול · סיכום AI · חיפוש · ניהול מטופלים — הכל אוטומטי",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# AI Pipeline
add_shape_rect(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.6), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.7), Inches(1.8), Inches(11.9), Inches(0.4),
             "🔄 AI Processing Pipeline (Celery Async)", font_size=16, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

pipeline = "🎤 Recording  →  ☁️ Upload to S3  →  🗣️ Whisper STT  →  🤖 Claude Summary  →  🏷️ Auto-Tagging  →  🔔 WebSocket Notify"
add_text_box(slide, Inches(0.7), Inches(2.3), Inches(11.9), Inches(0.7),
             pipeline, font_size=16, color=LIGHT, alignment=PP_ALIGN.CENTER)

# Backend features
add_shape_rect(slide, Inches(0.5), Inches(3.5), Inches(6.0), Inches(3.5), BG_CARD, GREEN)
add_text_box(slide, Inches(0.7), Inches(3.6), Inches(5.6), Inches(0.4),
             "⚙️ Backend — FastAPI 3.11+", font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

backend_items = [
    "60+ API Endpoints · 18 Routers · 21 Services",
    "16 DB Models — SQLAlchemy + Alembic Migrations",
    "Google OAuth → JWT (HS256/RS256) · Auto-refresh",
    "Celery Workers — STT + LLM + Tagging pipeline",
    "Meilisearch — Full-text + Semantic search",
    "MinIO/S3 — Audio storage with AES-256",
    "Structured logging (structlog) + Sentry",
    "Prometheus metrics + Rate limiting (slowapi)",
]
add_bullet_list(slide, Inches(0.7), Inches(4.1), Inches(5.6), Inches(2.8),
                backend_items, font_size=12, color=LIGHT, alignment=PP_ALIGN.LEFT)

# Frontend features
add_shape_rect(slide, Inches(6.8), Inches(3.5), Inches(6.0), Inches(3.5), BG_CARD, ORANGE)
add_text_box(slide, Inches(7.0), Inches(3.6), Inches(5.6), Inches(0.4),
             "📱 Frontend — Flutter Web + Dart 3", font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

frontend_items = [
    "15 Screens + 12 Riverpod Providers",
    "Audio Recording — MediaRecorder API בדפדפן",
    "Real-time updates via WebSocket",
    "8 שפות (he, en, de, es, fr, pt, ko, it)",
    "RTL/LTR dynamic switching",
    "Material 3 Design + Dark/Light mode",
    "Chunked upload with progress indicator",
    "Auto JWT refresh on 401 (Dio interceptor)",
]
add_bullet_list(slide, Inches(7.0), Inches(4.1), Inches(5.6), Inches(2.8),
                frontend_items, font_size=12, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 6 — Doctor Scribe AI — Data Flow
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "Doctor Scribe AI — Full Data Flow", font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

flow_steps = [
    ("🎤", "הקלטת ביקור", "רופא מקליט דרך\nMediaRecorder API\nבדפדפן", ACCENT),
    ("☁️", "העלאה ל-S3", "Chunked Upload\nMinIO Storage\nAES-256 Encrypted", GREEN),
    ("🗣️", "תמלול (STT)", "OpenAI Whisper API\nזיהוי דוברים\nעברית + אנגלית", ORANGE),
    ("🤖", "סיכום AI", "Claude LLM\nתלונה · ממצאים\nאבחנה · המלצות", PURPLE),
    ("🏷️", "תיוג אוטומטי", "אבחנות · תרופות\nתסמינים · ICD-10\nConfidence Score", RED),
    ("🔍", "אינדוקס + חיפוש", "Meilisearch\nFull-text + Semantic\nFilter by tags/date", ACCENT2),
]

step_w = Inches(1.85)
step_h = Inches(3.2)
x_start = Inches(0.4)
gap = Inches(0.23)
y = Inches(1.3)

flow_tooltips = [
    "רופא מקליט דרך MediaRecorder API בדפדפן, פורמט WebM/WAV",
    "Chunked Upload ל-MinIO, הצפנת AES-256, חתימת URL מאובטחת",
    "Whisper API: תמלול אודיו לטקסט, זיהוי דוברים, עברית+אנגלית",
    "Claude LLM: סיכום מובנה — תלונה, ממצאים, אבחנה, המלצות",
    "חילוץ אבחנות, תרופות, תסמינים, קודי ICD-10 עם Confidence",
    "Meilisearch: חיפוש Full-text + Semantic עם פילטרים מתקדמים",
]

for i, (icon, title, desc, color) in enumerate(flow_steps):
    x = x_start + i * (step_w + gap)
    add_shape_rect(slide, x, y, step_w, step_h, BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.15), step_w, Inches(0.5),
                 icon, font_size=32, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.7), step_w - Inches(0.2), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.2), step_w - Inches(0.2), Inches(1.6),
                 desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_tooltip_overlay(slide, x, y, step_w, step_h, flow_tooltips[i])
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + step_w, y + Inches(1.2), Inches(gap), Inches(0.5),
                     "→", font_size=22, color=GRAY, alignment=PP_ALIGN.CENTER)

# Business value
add_shape_rect(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.2), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.7), Inches(4.95), Inches(11.9), Inches(0.4),
             "💼 Business Impact", font_size=18, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

biz = [
    "▸ חוסך 15-20 דקות לרופא בכל ביקור — תיעוד אוטומטי במקום כתיבה ידנית",
    "▸ דיוק גבוה בסיכומים — LLM מחלץ אבחנות, תרופות והמלצות מהתמלול",
    "▸ חיפוש רפואי מתקדם — רופא מוצא ביקורים קודמים לפי אבחנה, תרופה או תסמין",
    "▸ Compliance מובנה — כל פעולה נרשמת ב-Audit Log, מידע מוצפן AES-256",
]
add_bullet_list(slide, Inches(0.7), Inches(5.45), Inches(11.9), Inches(1.5),
                biz, font_size=14, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 7 — Parent Website
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "🌐 Parent Website — drsscribe.com", font_size=36, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
             "מערכת ניהול תוכן רפואי · פורום קהילתי · כתיבת מאמרים אוטומטית",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Features — 3 columns
col_w = Inches(3.8)
col_h = Inches(3.0)
x_start = Inches(0.75)
gap = Inches(0.35)
y = Inches(1.7)

feature_cols = [
    ("📰 מערכת מאמרים", GREEN, [
        "150+ מאמרים רפואיים",
        "סינון לפי קטגוריה + תגיות",
        "מיון לפי צפיות / תאריך",
        "SEO Optimized (Next.js SSR)",
        "Responsive cards + sidebar",
        "כתיבת מאמרים AI (GPT-4.1)",
    ]),
    ("💬 פורום שאלות ותשובות", ACCENT, [
        "Stack Overflow Style Design",
        "Voting (Upvote/Downvote)",
        "תשובה מאושרת + ניקוד (Reputation)",
        "תגיות + חיפוש + מיון",
        "WYSIWYG Editor (Rich Text)",
        "JWT Auth + Auto-refresh",
    ]),
    ("🛠️ Tech Stack", PURPLE, [
        "Next.js 15 + React 18 + TypeScript",
        "Tailwind CSS + Custom Themes",
        "RTL Support (Hebrew-first)",
        "17+ Routes + Dynamic Pages",
        "FastAPI Backend Integration",
        "Docker Containerized Deploy",
    ]),
]

website_tooltips = [
    "150+ מאמרים רפואיים עם סינון, מיון, SEO, וכתיבת תוכן AI",
    "Stack Overflow style: הצבעות, תשובה מאושרת, Reputation, WYSIWYG",
    "Next.js 15, React 18, TypeScript, Tailwind CSS, RTL, Docker",
]

for i, (title, color, items) in enumerate(feature_cols):
    x = x_start + i * (col_w + gap)
    add_shape_rect(slide, x, y, col_w, col_h, BG_CARD, color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), col_w - Inches(0.3), Inches(0.4),
                 title, font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.2), y + Inches(0.55), col_w - Inches(0.4), Inches(2.3),
                    items, font_size=12, color=LIGHT, alignment=PP_ALIGN.LEFT)
    add_tooltip_overlay(slide, x, y, col_w, col_h, website_tooltips[i])

# Business value
add_shape_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), BG_CARD, GREEN)
add_text_box(slide, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.4),
             "💼 Business Value", font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

web_biz = [
    "▸ פלטפורמת שיווק — מושכת רופאים דרך תוכן רפואי איכותי (Content Marketing)",
    "▸ פורום קהילתי — בונה קהילה של רופאים סביב המוצר, מגביר Engagement ו-Retention",
    "▸ Lead Generation — רופאים שנרשמים לפורום הופכים ללקוחות פוטנציאליים ל-Doctor Scribe AI",
    "▸ SEO Organic — מאמרים מקצועיים מביאים תנועה אורגנית מגוגל",
]
add_bullet_list(slide, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.3),
                web_biz, font_size=13, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 8 — Security & Compliance
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), RED)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "🔒 Security & Compliance — Medical Grade", font_size=36, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.5),
             "אבטחה ברמה רפואית — הצפנה, RBAC, Audit Log, PII Masking",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Security cards — 2x3 grid
sec_cards = [
    ("🔐", "AES-256 Encryption", "כל נתוני המטופלים מוצפנים\nשם · ת.ז · טלפון · אלרגיות\nEncryption at rest + in transit", RED),
    ("🛡️", "RBAC + JWT Auth", "Google OAuth 2.0 → JWT\nHS256 / RS256 signing\ndoctor_id == current_user", ACCENT),
    ("📋", "Audit Logging", "כל פעולה נרשמת עם\nuser · action · timestamp · IP\nCompliance-ready", GREEN),
    ("🕵️", "PII Masking", "מסיכת מידע אישי אוטומטית\nRegex: ת.ז · טלפון · אימייל\nבלוגים ובתצוגות", ORANGE),
    ("⏱️", "Rate Limiting", "100 req/min API כללי\n10 req/min Authentication\nSliding window (slowapi)", PURPLE),
    ("🌐", "TLS 1.3 + HSTS", "SSL/TLS 1.3 enforced\nHSTS max-age=31536000\nCSP · X-Frame-Options · nosniff", ACCENT2),
]

card_w = Inches(3.8)
card_h = Inches(2.2)
x_start = Inches(0.75)
gap_x = Inches(0.35)
gap_y = Inches(0.25)
y_start = Inches(1.7)

security_tooltips = [
    "הצפנת כל נתוני המטופלים: שם, ת.ז, טלפון, אלרגיות. At rest + in transit",
    "Google OAuth 2.0 → JWT tokens, בדיקת הרשאות בכל endpoint",
    "כל פעולה נרשמת: משתמש, פעולה, זמן, IP. מוכן ל-Compliance",
    "מסיכת מידע אישי אוטומטית בלוגים: ת.ז, טלפון, אימייל → ***",
    "הגנה מ-DDoS: 100 req/min כללי, 10 req/min לאימות",
    "SSL/TLS 1.3 enforced, HSTS, CSP headers, X-Frame-Options",
]

for idx, (icon, title, desc, color) in enumerate(sec_cards):
    col = idx % 3
    row = idx // 3
    x = x_start + col * (card_w + gap_x)
    y = y_start + row * (card_h + gap_y)
    add_shape_rect(slide, x, y, card_w, card_h, BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.1), card_w, Inches(0.5),
                 f"{icon} {title}", font_size=15, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.65), card_w - Inches(0.3), Inches(1.4),
                 desc, font_size=12, color=LIGHT, alignment=PP_ALIGN.CENTER)
    add_tooltip_overlay(slide, x, y, card_w, card_h, security_tooltips[idx])


# ════════════════════════════════════════════════════
# SLIDE 9 — Challenges + Complexity
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ORANGE)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             "🧗 Challenges & Complexity", font_size=36, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

# Left — Challenges
add_shape_rect(slide, Inches(0.5), Inches(1.4), Inches(6.0), Inches(5.5), BG_CARD, ORANGE)
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.6), Inches(0.4),
             "⚠️ אתגרים טכניים שנפתרו", font_size=18, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

challenges = [
    "🎯 תמלול עברית — Whisper מתקשה עם עברית רפואית;\nפתרון: fine-tuning הפרומפט + post-processing",
    "🔄 סנכרון 9 שירותים — Docker health checks,\ndependency ordering, graceful restarts",
    "📦 Chunked Upload — קבצי אודיו גדולים (>100MB);\nפתרון: upload בחלקים עם resume capability",
    "🌍 RTL + LTR — תמיכה ב-8 שפות עם כיווניות דינמית;\n486 מפתחות תרגום × 8 שפות = 3,888 strings",
    "🔐 הצפנת PII — encrypt/decrypt שקוף לאפליקציה;\nAES-256 עם key rotation support",
    "⚡ Real-time — Celery + WebSocket לעדכון מיידי;\ntask chaining: STT → Summary → Tags → Notify",
]
add_bullet_list(slide, Inches(0.7), Inches(2.0), Inches(5.6), Inches(4.5),
                challenges, font_size=12, color=LIGHT, alignment=PP_ALIGN.LEFT)

# Right — Complexity metrics
add_shape_rect(slide, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.5), BG_CARD, PURPLE)
add_text_box(slide, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4),
             "📊 מורכבות הפרויקט", font_size=18, color=PURPLE, bold=True, alignment=PP_ALIGN.CENTER)

metrics = [
    "📁 ~50,000+ שורות קוד (Backend בלבד)",
    "🗂️ 3 שפות תכנות: Python · Dart · TypeScript",
    "🧩 9 Docker Microservices",
    "🔌 60+ API Endpoints",
    "🗄️ 16+ Database Models + Migrations",
    "📱 15 Flutter Screens + 12 Providers",
    "🌐 17+ Next.js Routes",
    "🤖 12+ ML Models",
    "🌍 8 Languages × 486 keys = 3,888 translations",
    "🔐 5 Security Layers (Auth→RBAC→Audit→Encrypt→Rate)",
    "📊 Prometheus + Grafana Monitoring",
    "🧪 Pytest + Integration Tests",
]
add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.6), Inches(4.5),
                metrics, font_size=13, color=LIGHT, alignment=PP_ALIGN.LEFT)


# ════════════════════════════════════════════════════
# SLIDE 10 — Summary & Thank You
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.5), Inches(0.8), Inches(12.3), Inches(1.0),
             "Summary — 3 Products, 1 Ecosystem", font_size=42, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Summary cards
summary = [
    ("🔬", "Data Science\nPlatform", "12+ ML Models\nModel Comparison\nToken-based SaaS", PURPLE),
    ("🏥", "Doctor Scribe\nAI", "60+ Endpoints\nAI Pipeline (STT→LLM)\nMedical-grade Security", ACCENT),
    ("🌐", "Parent\nWebsite", "150+ Articles\nSO-style Forum\nSEO + Content Marketing", GREEN),
]

card_w = Inches(3.6)
card_h = Inches(2.5)
gap = Inches(0.5)
x_start = Inches(1.3)
y = Inches(2.2)

summary_tooltips = [
    "פלטפורמת ML: 12+ מודלים, Token SaaS, דאשבורד Streamlit + API",
    "SaaS רפואי: 60+ endpoints, pipeline AI מלא, אבטחה ברמה רפואית",
    "150+ מאמרים, פורום קהילתי, SEO, כתיבת תוכן AI אוטומטית",
]

for i, (icon, title, desc, color) in enumerate(summary):
    x = x_start + i * (card_w + gap)
    add_shape_rect(slide, x, y, card_w, card_h, BG_CARD, color)
    add_text_box(slide, x, y + Inches(0.15), card_w, Inches(0.5),
                 icon, font_size=32, color=color, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.65), card_w - Inches(0.2), Inches(0.7),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), y + Inches(1.4), card_w - Inches(0.2), Inches(0.9),
                 desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)
    add_tooltip_overlay(slide, x, y, card_w, card_h, summary_tooltips[i])

# Skills demonstrated
add_shape_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.0), BG_CARD, ACCENT)
add_text_box(slide, Inches(0.7), Inches(5.15), Inches(11.9), Inches(0.8),
             "Python · FastAPI · Flutter · Next.js · React · TypeScript · Docker · PostgreSQL · Redis · "
             "scikit-learn · Whisper · Claude · Celery · Nginx · MinIO · Meilisearch",
             font_size=15, color=LIGHT, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8),
             "תודה רבה! 🙏",
             font_size=36, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)


# ── Save ──
output_path = r"C:\Doctor-Scribe\docs\MedScribe_AI_Portfolio.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
